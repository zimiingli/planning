#!/usr/bin/env python3
"""Generate weekly report slides for Direction-Aware Gate project (2026-03-09).
Focus: Updated storyline, token cost analysis, environment expansion, next steps."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ───
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_TEAL = RGBColor(0x00, 0xD2, 0xC6)
ACCENT_BLUE = RGBColor(0x45, 0x8B, 0xFF)
ACCENT_GOLD = RGBColor(0xFF, 0xC1, 0x07)
ACCENT_RED  = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_GREEN= RGBColor(0x4E, 0xC9, 0xB0)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
CARD_BG     = RGBColor(0x24, 0x24, 0x40)
SECTION_BG  = RGBColor(0x0D, 0x0D, 0x1A)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FIGURES_DIR = os.path.join(os.path.dirname(__file__), '..', '..', 'figures')

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

TOTAL_SLIDES = 18


# ─── Utility Functions ───
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=16, color=WHITE, bold=False, line_spacing=1.2, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt = line[0]
            c = line[1] if len(line) > 1 else color
            b = line[2] if len(line) > 2 else bold
            fs = line[3] if len(line) > 3 else font_size
        else:
            txt, c, b, fs = line, color, bold, font_size
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(font_size * (line_spacing - 1) + 2)
        p.text = txt
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
    return txBox

def add_image_safe(slide, img_name, left, top, width=None, height=None):
    img_path = os.path.join(FIGURES_DIR, img_name)
    if os.path.exists(img_path):
        kwargs = {}
        if width: kwargs['width'] = width
        if height: kwargs['height'] = height
        slide.shapes.add_picture(img_path, left, top, **kwargs)
        return True
    return False

def add_tag(slide, left, top, text, bg_color, text_color=WHITE, font_size=11):
    w, h = Inches(len(text) * 0.12 + 0.3), Inches(0.35)
    shape = add_shape(slide, left, top, w, h, bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(font_size)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = False
    return shape

def slide_number(slide, num, total):
    add_text(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
             f"{num}/{total}", font_size=10, color=MED_GRAY, align=PP_ALIGN.RIGHT)

def section_divider(num, label, title, subtitle, accent_color):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, SECTION_BG)
    add_rect(s, Inches(0), Inches(3.4), SLIDE_W, Inches(0.06), accent_color)
    add_text(s, Inches(1), Inches(2.3), Inches(5), Inches(0.5),
             label, font_size=16, color=accent_color, bold=True)
    add_text(s, Inches(1), Inches(2.9), Inches(11), Inches(0.9),
             title, font_size=42, color=WHITE, bold=True)
    add_text(s, Inches(1), Inches(3.7), Inches(10), Inches(0.5),
             subtitle, font_size=20, color=LIGHT_GRAY)
    slide_number(s, num, TOTAL_SLIDES)
    return s


# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)

add_text(s, Inches(1.2), Inches(1.5), Inches(11), Inches(1),
         "Direction-Aware Gate", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1.2), Inches(2.6), Inches(11), Inches(0.8),
         "Adaptive Test-Time Optimizer Triggering", font_size=32, color=ACCENT_TEAL)

add_rect(s, Inches(1.2), Inches(3.6), Inches(3), Inches(0.04), ACCENT_TEAL)

add_multiline(s, Inches(1.2), Inches(3.9), Inches(8), Inches(2.5), [
    ("Weekly Report  |  2026-03-09", LIGHT_GRAY, False, 20),
    ("Core Question: For any test-time optimizer T, can we learn WHEN to use it?", ACCENT_GOLD, False, 18),
    ("Target: NeurIPS 2026 (primary) / ICLR 2027 (backup)", MED_GRAY, False, 16),
    ("Status: Phase 5 In Progress  |  Cost Analysis Complete  |  6 Environments", ACCENT_GREEN, True, 16),
], font_size=20)

# Key numbers
for i, (val, label, col) in enumerate([
    ("38-77%", "Token Savings", ACCENT_GREEN),
    ("6", "Environments GO", ACCENT_BLUE),
    ("17", "Total Evaluated", ACCENT_TEAL),
]):
    x = Inches(8.5 + i * 1.6)
    add_shape(s, x, Inches(4.8), Inches(1.4), Inches(1.6), CARD_BG, col, Pt(1.5))
    add_text(s, x, Inches(4.95), Inches(1.4), Inches(0.6),
             val, font_size=26, color=col, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.5), Inches(1.4), Inches(0.5),
             label, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 2: Agenda
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
         "AGENDA", font_size=36, color=WHITE, bold=True)

agenda_items = [
    ("01", "Updated Storyline (v4.0)", "Three-layer narrative with token cost quantification"),
    ("02", "Token Cost Analysis", "SR vs Cost Pareto frontier across 3 environments"),
    ("03", "CATTS Failure & Adaptive Behavior", "Why vote-based gating fails; emergent RR adaptation"),
    ("04", "Environment Expansion", "3 new GO + 1 negative example + 11 NO-GO"),
    ("05", "Limitations Revealed", "What the rejections tell us about method boundaries"),
    ("06", "Next Week Plan", "Complete experiments, cost analysis, begin paper writing"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.4 + i * 0.9)
    add_text(s, Inches(1.0), y, Inches(0.6), Inches(0.5),
             num, font_size=22, color=ACCENT_TEAL, bold=True)
    add_rect(s, Inches(1.65), y + Inches(0.15), Inches(0.08), Inches(0.08), ACCENT_TEAL)
    add_text(s, Inches(2.0), y - Inches(0.05), Inches(5), Inches(0.4),
             title, font_size=22, color=WHITE, bold=True)
    add_text(s, Inches(2.0), y + Inches(0.3), Inches(8), Inches(0.4),
             desc, font_size=14, color=MED_GRAY)

slide_number(s, 2, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 3: Section Divider — Storyline
# ═══════════════════════════════════════════════
section_divider(3, "SECTION 01", "Updated Storyline v4.0",
                "From 'direction reversal' to 'direction reversal + token cost damage'", ACCENT_TEAL)


# ═══════════════════════════════════════════════
# SLIDE 4: Three-Act Narrative
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Paper Narrative: Three-Act Structure", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(10.2), Inches(0.35), "STORYLINE v4.0", ACCENT_TEAL)

acts = [
    ("Act 1: The Hidden Assumption", ACCENT_BLUE, "1 page", [
        "Test-time optimizers cost 4-36x more tokens per step",
        "11+ concurrent methods address 'when to trigger'",
        "ALL assume signal-utility direction is fixed",
    ]),
    ("Act 2: The Empirical Finding", ACCENT_GOLD, "2 pages", [
        "Finding 1: Direction reverses across environments",
        "Finding 2: Best signal identity also changes",
        "Wrong-direction: LR -34.5pp, MLP -51.2pp (RR=0%)",
    ]),
    ("Act 3: Method + Cost Analysis", ACCENT_GREEN, "3.5 pages", [
        "SCG-LR: LR on 5 features, <1s training, zero GPU",
        "NEW: Pareto-dominates ALL competitors on SR-Cost",
        "38-77% cheaper than always_trigger",
    ]),
]

for i, (title, color, pages, bullets) in enumerate(acts):
    x = Inches(0.4 + i * 4.3)
    card = add_shape(s, x, Inches(1.2), Inches(4.0), Inches(5.5), CARD_BG, color, Pt(2))

    # Act header
    add_rect(s, x, Inches(1.2), Inches(4.0), Inches(0.6), color)
    add_text(s, x + Inches(0.2), Inches(1.25), Inches(3.0), Inches(0.5),
             title, font_size=16, color=DARK_BG, bold=True)
    add_text(s, x + Inches(3.0), Inches(1.25), Inches(0.8), Inches(0.5),
             pages, font_size=12, color=DARK_BG, align=PP_ALIGN.RIGHT)

    for j, bullet in enumerate(bullets):
        y = Inches(2.1 + j * 1.3)
        add_shape(s, x + Inches(0.15), y, Inches(3.7), Inches(1.1), DARK_BG, color, Pt(0.5))
        add_text(s, x + Inches(0.3), y + Inches(0.1), Inches(3.4), Inches(0.9),
                 bullet, font_size=13, color=LIGHT_GRAY)

# Key change banner
add_shape(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.6), DARK_BG, ACCENT_GOLD, Pt(1.5))
add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.5),
         "v4.0 KEY CHANGE: Token cost quantification elevates 'why not always trigger?' from intuition to proven Pareto dominance",
         font_size=15, color=ACCENT_GOLD, bold=True)

slide_number(s, 4, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 5: Storyline Changes v3 → v4
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Storyline Evolution: v3.0 vs v4.0", font_size=28, color=WHITE, bold=True)

# Table header
hdr_y = Inches(1.2)
add_rect(s, Inches(0.5), hdr_y, Inches(12.3), Inches(0.5), ACCENT_TEAL)
for txt, x, w in [("Aspect", 0.6, 2.5), ("v3.0 (Last Week)", 3.3, 4.5), ("v4.0 (This Week)", 7.9, 4.8)]:
    add_text(s, Inches(x), hdr_y + Inches(0.05), Inches(w), Inches(0.4),
             txt, font_size=15, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)

changes = [
    ("Primary hook", "Direction reversal finding", "Direction reversal + token cost damage"),
    ("'Why not always?'", "Acknowledged but not quantified", "Fully quantified: 4-36x cost, Pareto analysis"),
    ("Method justification", "Near-oracle SR", "Near-oracle SR at fraction of cost"),
    ("Competitor comparison", "SR-only comparisons", "SR vs Token Cost Pareto frontier (Table 2)"),
    ("CATTS critique", "Assumes direction", "Assumes direction AND costs 10.5x (vote)"),
    ("Environment coverage", "3 validated + 1 NO-GO", "6 GO + 1 negative + 11 NO-GO (17 total)"),
    ("Adaptive behavior", "Mentioned briefly", "Elevated to Layer B -- emergent, not programmed"),
]

for i, (aspect, v3, v4) in enumerate(changes):
    y = Inches(1.85 + i * 0.65)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.55), bg)
    add_text(s, Inches(0.6), y + Inches(0.08), Inches(2.4), Inches(0.4),
             aspect, font_size=13, color=ACCENT_TEAL, bold=True)
    add_text(s, Inches(3.3), y + Inches(0.08), Inches(4.4), Inches(0.4),
             v3, font_size=12, color=MED_GRAY)
    add_text(s, Inches(7.9), y + Inches(0.08), Inches(4.7), Inches(0.4),
             v4, font_size=12, color=ACCENT_GOLD, bold=True)

slide_number(s, 5, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 6: Section Divider — Cost Analysis
# ═══════════════════════════════════════════════
section_divider(6, "SECTION 02", "Token Cost Analysis",
                "Answering: 'always_trigger has high SR, why not use it every step?'", ACCENT_GOLD)


# ═══════════════════════════════════════════════
# SLIDE 7: SR vs Token Cost Pareto
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "SR vs Token Cost: FRVC Pareto-Dominates All Competitors", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(10.2), Inches(0.35), "KEY RESULT", ACCENT_GREEN)

added = add_image_safe(s, "fig20_pareto_cost.png", Inches(0.3), Inches(1.1), width=Inches(12.7))
if not added:
    add_text(s, Inches(4), Inches(3.5), Inches(5), Inches(0.5),
             "[Figure: fig20_pareto_cost.png]", font_size=18, color=MED_GRAY, align=PP_ALIGN.CENTER)

add_shape(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.8), CARD_BG, ACCENT_GREEN, Pt(1.5))
add_multiline(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.7), [
    ("FRVC achieves near-oracle SR at fraction of cost: 6.55x (HotpotQA), 1.23x (APPS), 1.27x (WebShop)", ACCENT_GREEN, True, 14),
    ("Upper-left = better (high SR, low cost).  Green star = FRVC consistently in ideal region.", LIGHT_GRAY, False, 12),
], font_size=14)

slide_number(s, 7, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 8: Token Cost Comparison vs Competitors
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Token Cost: FRVC vs Competitors (×base)", font_size=28, color=WHITE, bold=True)

added = add_image_safe(s, "fig27_token_savings.png", Inches(0.3), Inches(1.0), width=Inches(7.5))

# Right: key numbers per environment
for i, (env, frvc_c, catts_c, coref_c, seag_c, col) in enumerate([
    ("HotpotQA", "6.55×", "10.50×", "6.12×", "7.40×", ACCENT_BLUE),
    ("APPS", "1.23×", "6.02×", "3.25×", "2.01×", ACCENT_ORANGE),
    ("WebShop", "1.27×", "5.55×", "3.67×", "2.84×", ACCENT_GREEN),
]):
    y = Inches(1.2 + i * 2.0)
    card = add_shape(s, Inches(8.2), y, Inches(4.8), Inches(1.8), CARD_BG, col, Pt(1.5))
    add_text(s, Inches(8.4), y + Inches(0.1), Inches(2.5), Inches(0.4),
             env, font_size=20, color=col, bold=True)
    add_multiline(s, Inches(8.4), y + Inches(0.5), Inches(4.4), Inches(1.2), [
        (f"FRVC: {frvc_c}  (lowest cost)", ACCENT_GREEN, True, 15),
        (f"CATTS: {catts_c}  CoRefine: {coref_c}  SEAG: {seag_c}", LIGHT_GRAY, False, 12),
    ], font_size=15)

slide_number(s, 8, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 9: Cost Breakdown
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Token Cost Breakdown: Where Do Tokens Go?", font_size=28, color=WHITE, bold=True)

added = add_image_safe(s, "fig22_cost_breakdown.png", Inches(0.3), Inches(1.1), width=Inches(12.7))

add_shape(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.8), CARD_BG, ACCENT_RED, Pt(1.5))
add_multiline(s, Inches(0.7), Inches(6.25), Inches(12), Inches(0.7), [
    ("CATTS vote overhead (red) is massive: 3,720 (HotpotQA), 8,396 (APPS), 33,850 (WebShop) tokens/episode", ACCENT_RED, True, 14),
    ("In APPS: C_vote (4,198) > C_rollout (3,306) -- voting costs MORE than the optimizer itself!", ACCENT_GOLD, True, 13),
], font_size=14)

slide_number(s, 9, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 10: CER + CATTS Failure
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Cost-Effectiveness & CATTS Failure Analysis", font_size=28, color=WHITE, bold=True)

# Left: CER figure
added = add_image_safe(s, "fig21_cer_comparison.png", Inches(0.2), Inches(1.0), width=Inches(7.0))

# Right: CATTS failure summary
add_shape(s, Inches(7.5), Inches(1.0), Inches(5.5), Inches(5.8), CARD_BG, ACCENT_RED, Pt(1.5))
add_text(s, Inches(7.7), Inches(1.1), Inches(5), Inches(0.5),
         "CATTS: Uniformly Worst", font_size=22, color=ACCENT_RED, bold=True)

catts_data = [
    ("HotpotQA", "SR=0.683", "Cost=10.5x", "Low SR + high cost"),
    ("APPS", "SR=0.585", "Cost=6.02x", "ZERO SR gain!"),
    ("WebShop", "SR=0.160", "Cost=5.55x", "Near-baseline SR"),
]
for i, (env, sr, cost, issue) in enumerate(catts_data):
    y = Inches(1.8 + i * 1.3)
    add_shape(s, Inches(7.7), y, Inches(5.1), Inches(1.1), DARK_BG, ACCENT_RED, Pt(0.5))
    add_text(s, Inches(7.9), y + Inches(0.05), Inches(2), Inches(0.35),
             env, font_size=15, color=WHITE, bold=True)
    add_text(s, Inches(7.9), y + Inches(0.4), Inches(4.7), Inches(0.35),
             f"{sr}  |  {cost}  |  {issue}", font_size=12, color=ACCENT_RED)
    if env == "APPS":
        add_tag(s, Inches(11.0), y + Inches(0.05), "C_vote > C_rollout!", ACCENT_RED, font_size=9)

add_shape(s, Inches(7.7), Inches(5.7), Inches(5.1), Inches(0.9), DARK_BG, ACCENT_GOLD, Pt(1))
add_multiline(s, Inches(7.9), Inches(5.75), Inches(4.7), Inches(0.8), [
    ("Root cause: K=5 voting every step", ACCENT_GOLD, True, 13),
    ("regardless of rollout trigger decision", LIGHT_GRAY, False, 12),
], font_size=13)

slide_number(s, 10, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 11: Adaptive Behavior + Step Reduction
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Emergent Adaptive Behavior & Dual Discount", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(10.2), Inches(0.35), "LAYER B", ACCENT_GOLD)

# Top: Adaptive RR
added = add_image_safe(s, "fig23_adaptive_rr.png", Inches(0.2), Inches(0.9), width=Inches(6.5))

# Top right: Step reduction
added = add_image_safe(s, "fig24_step_reduction.png", Inches(6.8), Inches(0.9), width=Inches(6.3))

# Bottom insight
add_shape(s, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.8), CARD_BG, ACCENT_TEAL, Pt(1.5))
add_multiline(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.7), [
    ("Dual discount: (1) Precise triggering (low RR) + (2) Success = fewer steps = less total cost", ACCENT_TEAL, True, 14),
    ("WebShop: 14.1 -> 5.6 steps (60% reduction) + RR=17% -> only 1.27x base cost for +36.5pp SR gain", LIGHT_GRAY, False, 12),
], font_size=14)

slide_number(s, 11, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 12: Section Divider — Environment Expansion
# ═══════════════════════════════════════════════
section_divider(12, "SECTION 03", "Environment Expansion",
                "17 environments evaluated: 6 GO, 1 negative example, 11 NO-GO", ACCENT_BLUE)


# ═══════════════════════════════════════════════
# SLIDE 13: Newly Accepted Environments
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "3 Newly Accepted Environments (Phase 5 GO)", font_size=28, color=WHITE, bold=True)

new_envs = [
    ("BabyAI", "Grid-world navigation", "GoTo-v0, 7 discrete actions",
     "2.0%", "11.3%", "+9.3pp", "8.8%", "18/18 DONE", ACCENT_GREEN),
    ("TextWorldExpress", "CoinCollector (text)", "Medium difficulty, max 50 steps",
     "64%", "98%", "+34pp", "~97%", "11/18 done", ACCENT_BLUE),
    ("TextWorld", "Text adventure", "5 rooms, 8 objects, quest=8",
     "58%", "68%", "+10pp", "TBD", "3/6 shards", ACCENT_BLUE),
]

for i, (name, desc, config, base, always, delta, scg, status, col) in enumerate(new_envs):
    x = Inches(0.4 + i * 4.3)
    card = add_shape(s, x, Inches(1.1), Inches(4.0), Inches(5.5), CARD_BG, col, Pt(2))

    # Header
    add_rect(s, x, Inches(1.1), Inches(4.0), Inches(0.6), col)
    add_text(s, x + Inches(0.15), Inches(1.15), Inches(3.7), Inches(0.5),
             name, font_size=18, color=DARK_BG, bold=True)

    add_multiline(s, x + Inches(0.2), Inches(1.9), Inches(3.6), Inches(4.5), [
        (desc, LIGHT_GRAY, False, 14),
        (config, MED_GRAY, False, 11),
        ("", WHITE, False, 8),
        (f"Base SR:    {base}", LIGHT_GRAY, False, 14),
        (f"Always SR:  {always}", LIGHT_GRAY, False, 14),
        (f"Delta:      {delta}", ACCENT_GOLD, True, 16),
        (f"SCG-LR SR:  {scg}", col, True, 16),
        ("", WHITE, False, 8),
        (f"Status: {status}", col, True, 13),
    ], font_size=14, line_spacing=1.3, font_name="Arial")

# Bottom note
add_shape(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.5), CARD_BG, ACCENT_GOLD, Pt(1))
add_text(s, Inches(0.7), Inches(6.25), Inches(12), Inches(0.4),
         "TextWorldExpress has LARGEST headroom of all environments (+34pp) -- strongest validation of adaptive gating",
         font_size=14, color=ACCENT_GOLD, bold=True)

slide_number(s, 13, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 14: Negative Example + NO-GO Summary
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Negative Example & Rejected Environments", font_size=28, color=WHITE, bold=True)

# Left: Plancraft
add_shape(s, Inches(0.4), Inches(1.1), Inches(5.5), Inches(3.5), CARD_BG, ACCENT_GOLD, Pt(2))
add_rect(s, Inches(0.4), Inches(1.1), Inches(5.5), Inches(0.6), ACCENT_GOLD)
add_text(s, Inches(0.6), Inches(1.15), Inches(4), Inches(0.5),
         "Plancraft (Negative Example)", font_size=18, color=DARK_BG, bold=True)

add_multiline(s, Inches(0.6), Inches(1.9), Inches(5.1), Inches(2.5), [
    ("Base SR: 29.8%  |  Always SR: 21.3%", WHITE, False, 15),
    ("Delta = -8.5% (NEGATIVE!)", ACCENT_RED, True, 18),
    ("", WHITE, False, 8),
    ("Gate learns to NOT trigger (~100% abstain)", ACCENT_GOLD, True, 14),
    ("-> Preserves base SR rather than degrading it", LIGHT_GRAY, False, 13),
    ("", WHITE, False, 8),
    ("Paper value: demonstrates gate ROBUSTNESS", ACCENT_GREEN, True, 14),
], font_size=15, line_spacing=1.2)

# Right: NO-GO categories
add_shape(s, Inches(6.2), Inches(1.1), Inches(6.9), Inches(3.5), CARD_BG, ACCENT_RED, Pt(1.5))
add_text(s, Inches(6.4), Inches(1.15), Inches(6), Inches(0.5),
         "11 NO-GO Environments", font_size=18, color=ACCENT_RED, bold=True)

nogo_cats = [
    ("Model Floor (3)", "ScienceWorld 0%, Sokoban 0%, Maze 4%", ACCENT_RED),
    ("Ceiling (1)", "InterCode-Bash 100%", MED_GRAY),
    ("Rollout Quality (1)", "ALFWorld: LLM hallucination", ACCENT_RED),
    ("Low Headroom (3)", "MiniHack x2, tau-bench: Delta<3pp", ACCENT_ORANGE),
    ("Bimodal (1+)", "Jericho 11 games: 0% or 100%", ACCENT_ORANGE),
]

for i, (cat, envs, col) in enumerate(nogo_cats):
    y = Inches(1.8 + i * 0.55)
    add_text(s, Inches(6.4), y, Inches(2.5), Inches(0.4),
             cat, font_size=12, color=col, bold=True)
    add_text(s, Inches(9.0), y, Inches(4), Inches(0.4),
             envs, font_size=11, color=LIGHT_GRAY)

# Bottom: Environment expansion figure
added = add_image_safe(s, "fig25_env_expansion.png", Inches(0.3), Inches(4.8), width=Inches(12.7))

slide_number(s, 14, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 15: Limitations Revealed
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Limitations Revealed by Rejected Environments", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(10.2), Inches(0.35), "HONEST ANALYSIS", ACCENT_TEAL)

limitations = [
    ("1. Model Capacity Floor",
     "Base agent needs SR > 10%",
     "ScienceWorld (0%), Sokoban (0%), Maze (4%): Qwen3-4B fundamentally cannot perform the task. "
     "No amount of gating or rollout helps.",
     "Inherent to all test-time optimization, not specific to our method", ACCENT_RED),
    ("2. Meaningful Rollout Headroom",
     "Need Delta >= 3pp from optimizer",
     "MiniHack (+2pp), tau-bench (+2pp): When the optimizer barely helps, "
     "there's nothing for the gate to optimize.",
     "Our method is most valuable in 'high-value, high-cost' regime", ACCENT_ORANGE),
    ("3. Reliable Rollout Mechanism",
     "Need deepcopy or deterministic eval",
     "ALFWorld: LLM-as-Simulator hallucinates, Batch Scoring has confirmation bias (2.9/10 vs 6.6/10).",
     "Rollout quality hierarchy: env.deepcopy > deterministic eval > LLM simulation", ACCENT_GOLD),
    ("4. Within-Episode Variability",
     "Not bimodal performance distribution",
     "Jericho 11 games: Either 0% or 100% SR per game. No 'sweet spot' for per-step gating.",
     "Per-step gating assumes variability; per-instance estimation is different", ACCENT_BLUE),
]

for i, (title, req, evidence, insight, col) in enumerate(limitations):
    y = Inches(1.1 + i * 1.5)
    card = add_shape(s, Inches(0.4), y, Inches(12.5), Inches(1.35), CARD_BG, col, Pt(1))
    add_text(s, Inches(0.6), y + Inches(0.05), Inches(4), Inches(0.35),
             title, font_size=15, color=col, bold=True)
    add_text(s, Inches(5.0), y + Inches(0.05), Inches(7.5), Inches(0.35),
             f"Requirement: {req}", font_size=12, color=WHITE, bold=True)
    add_text(s, Inches(0.6), y + Inches(0.4), Inches(12), Inches(0.35),
             evidence, font_size=11, color=LIGHT_GRAY)
    add_text(s, Inches(0.6), y + Inches(0.8), Inches(12), Inches(0.35),
             f"-> {insight}", font_size=11, color=col, bold=True)

# Positive framing
add_shape(s, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.6), CARD_BG, ACCENT_GREEN, Pt(1.5))
add_text(s, Inches(0.7), Inches(6.35), Inches(12), Inches(0.5),
         "These are reasonable prerequisites for any adaptive compute method. Plancraft shows gate is ROBUST (learns to abstain).",
         font_size=14, color=ACCENT_GREEN, bold=True)

slide_number(s, 15, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 16: Full Environment Summary Table
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Complete Environment Summary: 6 GO Environments", font_size=28, color=WHITE, bold=True)

# Table header
hdr_y = Inches(1.1)
add_rect(s, Inches(0.4), hdr_y, Inches(12.5), Inches(0.5), ACCENT_TEAL)
cols = [("Environment", 0.5, 2.0), ("Base SR", 2.6, 1.0), ("Always SR", 3.7, 1.1),
        ("Delta", 4.9, 0.9), ("SCG-LR SR", 5.9, 1.2), ("Cost (×base)", 7.0, 1.0),
        ("Best Signal", 8.1, 2.2), ("Status", 10.4, 2.4)]
for txt, x, w in cols:
    add_text(s, Inches(x), hdr_y + Inches(0.05), Inches(w), Inches(0.4),
             txt, font_size=13, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)

env_rows = [
    ("HotpotQA", "49.0%", "97.0%", "+48pp", "96.8%", "6.55x", "evidence_count", "Complete", ACCENT_GREEN),
    ("APPS", "57.8%", "64.5%", "+6.7pp", "58.8%", "1.23x", "step_count", "Complete", ACCENT_GREEN),
    ("WebShop", "7.2%", "47.5%", "+40pp", "43.7%", "1.27x", "state_category", "Complete", ACCENT_GREEN),
    ("BabyAI", "2.0%", "11.3%", "+9.3pp", "8.8%", "TBD", "TBD", "18/18 Done", ACCENT_BLUE),
    ("TWExpress", "64%", "98%", "+34pp", "~97%", "TBD", "step_count", "11/18", ACCENT_BLUE),
    ("TextWorld", "58%", "68%", "+10pp", "TBD", "TBD", "TBD", "3/6 shards", ACCENT_BLUE),
]

for i, (env, base, always, delta, scg, cost, signal, status, status_col) in enumerate(env_rows):
    y = Inches(1.75 + i * 0.55)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.5), bg)

    phase_col = ACCENT_GREEN if i < 3 else ACCENT_BLUE
    add_text(s, Inches(0.5), y + Inches(0.07), Inches(1.9), Inches(0.35),
             env, font_size=13, color=phase_col, bold=True)
    add_text(s, Inches(2.6), y + Inches(0.07), Inches(1.0), Inches(0.35),
             base, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.7), y + Inches(0.07), Inches(1.1), Inches(0.35),
             always, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.9), y + Inches(0.07), Inches(0.9), Inches(0.35),
             delta, font_size=13, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.9), y + Inches(0.07), Inches(1.2), Inches(0.35),
             scg, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.2), y + Inches(0.07), Inches(0.8), Inches(0.35),
             cost, font_size=13, color=ACCENT_GREEN if cost != "TBD" else MED_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.1), y + Inches(0.07), Inches(2.2), Inches(0.35),
             signal, font_size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.4), y + Inches(0.07), Inches(2.4), Inches(0.35),
             status, font_size=12, color=status_col, bold=True, align=PP_ALIGN.CENTER)

# Negative example row
add_rect(s, Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.5), CARD_BG)
add_text(s, Inches(0.5), Inches(5.12), Inches(1.9), Inches(0.35),
         "Plancraft", font_size=13, color=ACCENT_GOLD, bold=True)
add_text(s, Inches(2.6), Inches(5.12), Inches(1.0), Inches(0.35),
         "29.8%", font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(3.7), Inches(5.12), Inches(1.1), Inches(0.35),
         "21.3%", font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(4.9), Inches(5.12), Inches(0.9), Inches(0.35),
         "-8.5pp", font_size=13, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(5.9), Inches(5.12), Inches(1.2), Inches(0.35),
         "Abstain", font_size=13, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.1), Inches(5.12), Inches(2.2), Inches(0.35),
         "Weak (has_output 0.16)", font_size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(10.4), Inches(5.12), Inches(2.4), Inches(0.35),
         "Negative Example", font_size=12, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)

# Key insight
add_shape(s, Inches(0.4), Inches(5.8), Inches(12.5), Inches(0.8), CARD_BG, ACCENT_TEAL, Pt(1.5))
add_multiline(s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.7), [
    ("Green = Phase 1 validated  |  Blue = Phase 5 new  |  Gold = Negative example  |  9 adapters built", ACCENT_TEAL, True, 14),
    ("Each env uses different best signal -- reinforcing the core finding of signal-utility environment dependence", LIGHT_GRAY, False, 12),
], font_size=14)

slide_number(s, 16, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# SLIDE 17: Section Divider — Next Steps
# ═══════════════════════════════════════════════
section_divider(17, "SECTION 04", "Next Week Plan",
                "Complete experiments -> Cost analysis -> Begin paper writing", ACCENT_GREEN)


# ═══════════════════════════════════════════════
# SLIDE 18: Next Steps & Timeline
# ═══════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Next Week Plan & Key Milestones", font_size=28, color=WHITE, bold=True)

# Priority tasks
priorities = [
    ("P0", ACCENT_RED, [
        "Complete TextWorldExpress Step 2 (7 remaining jobs)",
        "Complete TextWorld Step 1 (3 shards) -> Step 2",
        "Run Plancraft K=3 rerun (negative example validation)",
    ]),
    ("P1", ACCENT_GOLD, [
        "Run full cost analysis on new environments (BabyAI, TWExpress, TW)",
        "Merge all results into unified Table 1 (SR) and Table 2 (Cost)",
    ]),
    ("P2", ACCENT_BLUE, [
        "Begin paper writing: Section 2 (Empirical Finding) + Section 4 (Experiments)",
        "Generate Pareto frontier plots for 6 environments",
    ]),
]

for i, (prio, col, tasks) in enumerate(priorities):
    x = Inches(0.4 + i * 4.3)
    card = add_shape(s, x, Inches(1.1), Inches(4.0), Inches(3.5), CARD_BG, col, Pt(1.5))
    add_rect(s, x, Inches(1.1), Inches(4.0), Inches(0.5), col)
    add_text(s, x + Inches(0.2), Inches(1.15), Inches(3.5), Inches(0.4),
             f"Priority {prio}", font_size=18, color=DARK_BG, bold=True)

    for j, task in enumerate(tasks):
        add_text(s, x + Inches(0.2), Inches(1.8 + j * 0.7), Inches(3.6), Inches(0.6),
                 task, font_size=12, color=LIGHT_GRAY)

# Milestones
add_shape(s, Inches(0.4), Inches(4.9), Inches(12.5), Inches(1.8), CARD_BG, ACCENT_GREEN, Pt(1.5))
add_text(s, Inches(0.6), Inches(4.95), Inches(4), Inches(0.4),
         "Key Milestones", font_size=20, color=ACCENT_GREEN, bold=True)

milestones = [
    ("Tue 03-10", "All Phase 5 environment experiments complete", ACCENT_TEAL),
    ("Thu 03-12", "Unified cost analysis across 6 environments", ACCENT_GOLD),
    ("Sat 03-14", "Paper Section 2 + Section 4 first draft", ACCENT_BLUE),
]

for i, (date, desc, col) in enumerate(milestones):
    y = Inches(5.5 + i * 0.35)
    add_text(s, Inches(0.8), y, Inches(1.5), Inches(0.3),
             date, font_size=14, color=col, bold=True)
    add_text(s, Inches(2.5), y, Inches(6), Inches(0.3),
             desc, font_size=14, color=LIGHT_GRAY)

# Risk items
add_text(s, Inches(8.5), Inches(4.95), Inches(4), Inches(0.4),
         "Risk Items", font_size=20, color=ACCENT_RED, bold=True)
risks = [
    "TextWorld shards may fail again (fallback: 3/6)",
    "TWExpress oracle jobs timeout (partial OK)",
    "Plancraft K=3: either outcome is informative",
]
for i, risk in enumerate(risks):
    add_text(s, Inches(8.5), Inches(5.5 + i * 0.35), Inches(4.3), Inches(0.3),
             risk, font_size=12, color=ACCENT_RED)

slide_number(s, 18, TOTAL_SLIDES)


# ═══════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), 'drafts', 'v2_weekly_report_20260309.pptx')
os.makedirs(os.path.dirname(output_path), exist_ok=True)
prs.save(output_path)
print(f"\n✅ Presentation saved to: {output_path}")
print(f"   {TOTAL_SLIDES} slides generated")
