#!/usr/bin/env python3
"""Generate weekly report presentation for Direction-Aware Gate project (2026-03-02)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ───
DARK_BG     = RGBColor(0x1A, 0x1A, 0x2E)  # Deep navy
ACCENT_TEAL = RGBColor(0x00, 0xD2, 0xC6)  # Bright teal
ACCENT_BLUE = RGBColor(0x45, 0x8B, 0xFF)  # Electric blue
ACCENT_GOLD = RGBColor(0xFF, 0xC1, 0x07)  # Gold
ACCENT_RED  = RGBColor(0xFF, 0x6B, 0x6B)  # Coral red
ACCENT_GREEN= RGBColor(0x4E, 0xC9, 0xB0)  # Mint green
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
CARD_BG     = RGBColor(0x24, 0x24, 0x40)  # Slightly lighter navy for cards
SECTION_BG  = RGBColor(0x0D, 0x0D, 0x1A)  # Darker for section dividers

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FIGURES_DIR = os.path.join(os.path.dirname(__file__), '..', '..', 'figures')

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ─── Utility Functions ───

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=16, color=WHITE, bold=False, line_spacing=1.2, bullet=False, font_name="Arial"):
    """Add multiple lines of text. lines can be list of str or list of (str, color, bold, font_size)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            txt, c, b = line[0], line[1] if len(line) > 1 else color, line[2] if len(line) > 2 else bold
            fs = line[3] if len(line) > 3 else font_size
        else:
            txt, c, b, fs = line, color, bold, font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(font_size * (line_spacing - 1) + 2)
        if bullet:
            p.text = f"  {txt}"
        else:
            p.text = txt
        p.font.size = Pt(fs)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = font_name
    return txBox

def add_image_safe(slide, img_name, left, top, width=None, height=None):
    img_path = os.path.join(FIGURES_DIR, img_name)
    if os.path.exists(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(img_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(img_path, left, top, height=height)
        else:
            slide.shapes.add_picture(img_path, left, top)
        return True
    return False

def add_tag(slide, left, top, text, bg_color, text_color=WHITE, font_size=11):
    w, h = Inches(len(text) * 0.12 + 0.3), Inches(0.35)
    shape = add_shape(slide, left, top, w, h, bg_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(font_size)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = "Arial"
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.word_wrap = False
    return shape

def slide_number(slide, num, total):
    add_text(slide, Inches(12.3), Inches(7.05), Inches(0.9), Inches(0.35),
             f"{num}/{total}", font_size=10, color=MED_GRAY, align=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 22

# ─────────────────────────────────────────────────
# SLIDE 1: Title
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(s, DARK_BG)

# Decorative accent line
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_TEAL)

# Title block
add_text(s, Inches(1.2), Inches(1.8), Inches(11), Inches(1),
         "Direction-Aware Gate", font_size=48, color=WHITE, bold=True)
add_text(s, Inches(1.2), Inches(2.9), Inches(11), Inches(0.8),
         "Adaptive Test-Time Optimizer Triggering", font_size=32, color=ACCENT_TEAL, bold=False)

# Divider
add_rect(s, Inches(1.2), Inches(3.9), Inches(3), Inches(0.04), ACCENT_TEAL)

# Subtitle info
add_multiline(s, Inches(1.2), Inches(4.2), Inches(8), Inches(2), [
    ("Weekly Report  |  2026-03-02", LIGHT_GRAY, False, 20),
    ("Core Question: For any test-time optimizer T, can we learn WHEN to use it?", ACCENT_GOLD, False, 18),
    ("Target: NeurIPS 2026 (primary) / ICLR 2027 (backup)", MED_GRAY, False, 16),
    ("Status: Phase 0-4 Complete  |  Phase 5 Ready to Launch", ACCENT_GREEN, True, 16),
], font_size=20)

slide_number(s, 1, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 2: Agenda
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
         "AGENDA", font_size=36, color=WHITE, bold=True)

agenda_items = [
    ("01", "Experiment Overview", "6 environments, base vs always, gate comparison"),
    ("02", "Finding 1: Direction Reversal", "Signal-utility direction is NOT fixed across environments"),
    ("03", "Finding 2: Signal Replacement", "Strongest signal identity differs per environment"),
    ("04", "Wrong-Direction Ablation", "Quantified catastrophic cost of wrong direction"),
    ("05", "Phase 5 Plan", "Dual-track: Cascaded Gate + ICGNet"),
    ("06", "Paper Storyline", "Three-act narrative + contribution ranking"),
    ("07", "Next Steps & Timeline", "Implementation roadmap"),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.4 + i * 0.8)
    # Number circle
    add_text(s, Inches(1.0), y, Inches(0.6), Inches(0.5),
             num, font_size=22, color=ACCENT_TEAL, bold=True, font_name="Arial")
    # Dot
    add_rect(s, Inches(1.65), y + Inches(0.15), Inches(0.08), Inches(0.08), ACCENT_TEAL)
    # Title
    add_text(s, Inches(2.0), y - Inches(0.05), Inches(5), Inches(0.4),
             title, font_size=22, color=WHITE, bold=True)
    # Description
    add_text(s, Inches(2.0), y + Inches(0.3), Inches(8), Inches(0.4),
             desc, font_size=14, color=MED_GRAY)

slide_number(s, 2, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 3: Section Divider — Experiment Overview
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, SECTION_BG)
add_rect(s, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), ACCENT_TEAL)
add_text(s, Inches(1), Inches(2.3), Inches(5), Inches(0.5),
         "SECTION 01", font_size=16, color=ACCENT_TEAL, bold=True)
add_text(s, Inches(1), Inches(2.9), Inches(11), Inches(0.9),
         "Full Experiment Data Overview", font_size=42, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.7), Inches(10), Inches(0.5),
         "6 Environments  ×  10 Methods  ×  3 Seeds  ×  200 Episodes", font_size=20, color=LIGHT_GRAY)
slide_number(s, 3, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 4: Environment Overview
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
         "Environment Overview & Optimizer Design", font_size=28, color=WHITE, bold=True)

# Table-like cards
envs = [
    ("HotpotQA", "Qwen3-4B", "K≤5 (search/lookup/finish)", "Per-action evaluation", "Main Exp", ACCENT_GREEN),
    ("APPS (Intro)", "Qwen3-4B", "Single (code gen)", "K-variant gen (K=5)", "2nd Effective", ACCENT_BLUE),
    ("WebShop", "Qwen3-4B", "Medium (search/click/buy)", "LLM-Propose-K (K=5,H=3)", "3rd Effective", ACCENT_BLUE),
    ("ALFWorld", "Qwen3-8B", "Medium discrete", "LLM-as-Sim / Batch Score", "NO-GO", ACCENT_RED),
    ("HumanEval", "Qwen3-4B", "Single (code gen)", "K-variant generation", "Ceiling", MED_GRAY),
    ("MBPP", "Qwen3-4B", "Single (code gen)", "K-variant generation", "Ceiling", MED_GRAY),
]

for i, (name, model, action, opt, role, role_color) in enumerate(envs):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.2 + row * 3.0)

    card = add_shape(s, x, y, Inches(3.9), Inches(2.7), CARD_BG, ACCENT_TEAL if "Main" in role else None, Pt(1.5) if "Main" in role else Pt(0))

    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(2.5), Inches(0.4),
             name, font_size=22, color=WHITE, bold=True)

    # Role tag
    tag_w = Inches(len(role) * 0.1 + 0.35)
    tag = add_shape(s, x + Inches(3.9) - tag_w - Inches(0.15), y + Inches(0.15), tag_w, Inches(0.32), role_color)
    tag.text_frame.paragraphs[0].text = role
    tag.text_frame.paragraphs[0].font.size = Pt(10)
    tag.text_frame.paragraphs[0].font.color.rgb = WHITE
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_multiline(s, x + Inches(0.2), y + Inches(0.65), Inches(3.5), Inches(1.8), [
        (f"Model: {model}", LIGHT_GRAY, False, 13),
        (f"Actions: {action}", LIGHT_GRAY, False, 13),
        (f"Optimizer T: {opt}", ACCENT_TEAL, False, 13),
    ], font_size=13, line_spacing=1.5)

slide_number(s, 4, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 5: Base vs Always-Trigger
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Base vs Always-Trigger: Does the Optimizer Help?", font_size=28, color=WHITE, bold=True)

# Data rows
data = [
    ("HotpotQA", "49.0%", "97.0%", "+48.0pp", "GO", ACCENT_GREEN),
    ("APPS (Intro)", "57.8%", "64.8%", "+7.0pp", "GO", ACCENT_GREEN),
    ("WebShop", "7.2%", "43.0%", "+35.8pp", "GO", ACCENT_GREEN),
    ("ALFWorld v3", "30.0%", "20.0%", "−10.0pp", "NO-GO", ACCENT_RED),
    ("HumanEval", "92.1%", "92.3%", "+0.2pp", "Ceiling", MED_GRAY),
    ("MBPP", "92.7%", "92.7%", "+0.0pp", "Ceiling", MED_GRAY),
]

# Header
hdr_y = Inches(1.2)
add_rect(s, Inches(0.5), hdr_y, Inches(12.3), Inches(0.5), ACCENT_TEAL)
cols = [("Environment", 0.6, 2.5), ("Base SR", 3.2, 1.5), ("Always SR", 4.8, 1.5),
        ("Δ (lift)", 6.4, 1.5), ("Verdict", 8.0, 1.5)]
for txt, x, w in cols:
    add_text(s, Inches(x), hdr_y + Inches(0.05), Inches(w), Inches(0.4),
             txt, font_size=16, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)

for i, (env, base, always, delta, verdict, vcolor) in enumerate(data):
    y = Inches(1.85 + i * 0.7)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.6), bg)

    add_text(s, Inches(0.7), y + Inches(0.1), Inches(2.3), Inches(0.4),
             env, font_size=16, color=WHITE, bold=True)
    add_text(s, Inches(3.2), y + Inches(0.1), Inches(1.5), Inches(0.4),
             base, font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.8), y + Inches(0.1), Inches(1.5), Inches(0.4),
             always, font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.4), y + Inches(0.1), Inches(1.5), Inches(0.4),
             delta, font_size=16, color=ACCENT_GOLD if "+" in delta and float(delta.replace("+","").replace("pp","")) > 5 else ACCENT_RED if "−" in delta else MED_GRAY,
             bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.0), y + Inches(0.1), Inches(1.5), Inches(0.4),
             verdict, font_size=16, color=vcolor, bold=True, align=PP_ALIGN.CENTER)

# Key insight box
add_shape(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9), CARD_BG, ACCENT_GOLD, Pt(1.5))
add_text(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.7),
         "Key: 3 effective environments (GO) + 1 boundary result (NO-GO) + 2 ceiling environments",
         font_size=16, color=ACCENT_GOLD, bold=True)

slide_number(s, 5, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 6: HotpotQA Main Results
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
         "HotpotQA: Main Experiment (10 Methods × 3 Seeds)", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(10.5), Inches(0.35), "MAIN EXPERIMENT", ACCENT_GREEN)

hotpot_data = [
    ("base_only", "49.0±1.9%", "0.0%", "100.0%", "0.000", WHITE),
    ("always_trigger", "97.0±0.4%", "100.0%", "0.0%", "0.000", WHITE),
    ("random_50", "89.0±0.8%", "51.4%", "48.6%", "0.614", WHITE),
    ("entropy_threshold", "67.2±3.3%", "21.5%", "78.5%", "0.509", WHITE),
    ("best_sigma_wrong", "58.2±2.5%", "49.9%", "50.1%", "0.277", ACCENT_RED),
    ("scg_prompt", "95.7±0.5%", "82.1%", "17.9%", "—", LIGHT_GRAY),
    ("scg_mlp", "96.7±0.6%", "63.7%", "36.3%", "—", LIGHT_GRAY),
    ("scg_finetune_lr ⭐", "96.7±0.6%", "55.9%", "44.1%", "0.609", ACCENT_GOLD),
    ("oracle", "97.0±0.4%", "33.0%", "67.0%", "0.802", ACCENT_TEAL),
]

hdr_y = Inches(1.1)
add_rect(s, Inches(0.4), hdr_y, Inches(12.5), Inches(0.45), ACCENT_TEAL)
hcols = [("Method", 0.5, 2.3), ("SR (mean±std)", 2.9, 1.8), ("RR (%)", 4.8, 1.2),
         ("CS (%)", 6.1, 1.2), ("TES", 7.4, 1.0)]
for txt, x, w in hcols:
    add_text(s, Inches(x), hdr_y + Inches(0.02), Inches(w), Inches(0.4),
             txt, font_size=14, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)

for i, (method, sr, rr, cs, tes, mcolor) in enumerate(hotpot_data):
    y = Inches(1.65 + i * 0.5)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.45), bg)

    add_text(s, Inches(0.5), y + Inches(0.05), Inches(2.2), Inches(0.35),
             method, font_size=13, color=mcolor, bold="⭐" in method or "wrong" in method)
    add_text(s, Inches(2.9), y + Inches(0.05), Inches(1.8), Inches(0.35),
             sr, font_size=13, color=mcolor, align=PP_ALIGN.CENTER)
    add_text(s, Inches(4.8), y + Inches(0.05), Inches(1.2), Inches(0.35),
             rr, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(6.1), y + Inches(0.05), Inches(1.2), Inches(0.35),
             cs, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, Inches(7.4), y + Inches(0.05), Inches(1.0), Inches(0.35),
             tes, font_size=13, color=mcolor, align=PP_ALIGN.CENTER, bold="⭐" in method)

# Key insight
add_shape(s, Inches(0.4), Inches(6.3), Inches(8.3), Inches(1.0), CARD_BG, ACCENT_GOLD, Pt(1))
add_multiline(s, Inches(0.6), Inches(6.4), Inches(7.8), Inches(0.8), [
    ("SCG-LR Pareto-dominates random: SR 96.7% >> 89.0%, comparable CS", ACCENT_GOLD, True, 14),
    ("Wrong-direction SR crashes to 58.2% (−38.8pp vs correct)", ACCENT_RED, True, 14),
], font_size=14)

# Figure
add_image_safe(s, "fig19_sr_comparison.png", Inches(9.0), Inches(1.6), height=Inches(5.0))

slide_number(s, 6, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 7: APPS + WebShop Results
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "APPS & WebShop: Supporting Environments", font_size=28, color=WHITE, bold=True)

# Left: APPS
add_shape(s, Inches(0.4), Inches(1.1), Inches(6.2), Inches(5.8), CARD_BG, ACCENT_BLUE, Pt(1))
add_text(s, Inches(0.6), Inches(1.2), Inches(4), Inches(0.5),
         "APPS Introductory", font_size=22, color=ACCENT_BLUE, bold=True)
add_tag(s, Inches(4.3), Inches(1.25), "2nd ENV", ACCENT_BLUE)

apps_lines = [
    ("base_only        SR 57.8%    RR 0%     TES 0.000", LIGHT_GRAY, False, 13),
    ("always_trigger   SR 64.8%    RR 100%   TES 0.000", LIGHT_GRAY, False, 13),
    ("random_50        SR 66.5%    RR 50.2%  TES 0.665", WHITE, False, 13),
    ("sigma_wrong      SR 58.5%    RR 0%     TES 0.174", ACCENT_RED, False, 13),
    ("scg_finetune_lr  SR 65.0%    RR 40.2%  TES 0.748", ACCENT_GOLD, True, 13),
    ("oracle           SR 66.8%    RR 100%   TES 0.000", ACCENT_TEAL, False, 13),
]
add_multiline(s, Inches(0.6), Inches(1.8), Inches(5.8), Inches(3.0), apps_lines, font_size=13, line_spacing=1.6, font_name="Courier New")

add_shape(s, Inches(0.6), Inches(4.6), Inches(5.8), Inches(0.9), DARK_BG, ACCENT_GOLD, Pt(1))
add_multiline(s, Inches(0.8), Inches(4.7), Inches(5.2), Inches(0.7), [
    ("TES_LR (0.748) > TES_random (0.665), p=0.001", ACCENT_GOLD, True, 13),
    ("Wrong-dir: RR=0% — passive abandonment mode", ACCENT_RED, True, 13),
], font_size=13)

# Right: WebShop
add_shape(s, Inches(6.9), Inches(1.1), Inches(6.2), Inches(5.8), CARD_BG, ACCENT_BLUE, Pt(1))
add_text(s, Inches(7.1), Inches(1.2), Inches(4), Inches(0.5),
         "WebShop", font_size=22, color=ACCENT_BLUE, bold=True)
add_tag(s, Inches(10.3), Inches(1.25), "3rd ENV", ACCENT_BLUE)

web_lines = [
    ("base_only        SR  7.2%   RR  0%   Prec —", LIGHT_GRAY, False, 13),
    ("always_trigger   SR 43.0%   RR 100%  Prec 12.9%", LIGHT_GRAY, False, 13),
    ("random_50        SR 47.5%   RR 50.9% Prec 21.9%", WHITE, False, 13),
    ("sigma_wrong      SR  7.2%   RR 37.1% Prec 0.0%", ACCENT_RED, False, 13),
    ("scg_finetune_lr  SR 43.7%   RR 16.9% Prec 75.1%", ACCENT_GOLD, True, 13),
    ("scg_lora         SR 42.8%   RR 17.7% Prec 72.4%", LIGHT_GRAY, False, 13),
    ("oracle           SR 43.3%   RR 13.1% Prec 100%", ACCENT_TEAL, False, 13),
]
add_multiline(s, Inches(7.1), Inches(1.8), Inches(5.8), Inches(3.5), web_lines, font_size=13, line_spacing=1.5, font_name="Courier New")

add_shape(s, Inches(7.1), Inches(4.6), Inches(5.8), Inches(0.9), DARK_BG, ACCENT_GOLD, Pt(1))
add_multiline(s, Inches(7.3), Inches(4.7), Inches(5.2), Inches(0.7), [
    ("SR ≈ oracle (43.7% vs 43.3%), precision 75.1%", ACCENT_GOLD, True, 13),
    ("6× compute efficiency (RR=16.9%) — best case!", ACCENT_GREEN, True, 13),
], font_size=13)

# Bottom: ALFWorld small note
add_shape(s, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.5), CARD_BG, ACCENT_RED, Pt(1))
add_text(s, Inches(0.7), Inches(6.2), Inches(11.8), Inches(0.4),
         "ALFWorld (NO-GO): Always-trigger SR 20% < Base 30%. Failure: confirmation bias, LLM hallucination scoring",
         font_size=14, color=ACCENT_RED, bold=True)

slide_number(s, 7, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 8: Pareto Front
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "SR vs Cost Saving: Pareto Front Across 3 Environments", font_size=28, color=WHITE, bold=True)

added = add_image_safe(s, "fig16_pareto_3env.png", Inches(0.8), Inches(1.2), width=Inches(11.5))
if not added:
    add_text(s, Inches(3), Inches(3.5), Inches(7), Inches(0.5),
             "[Figure: fig16_pareto_3env.png]", font_size=20, color=MED_GRAY, align=PP_ALIGN.CENTER)

add_shape(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.7), CARD_BG, ACCENT_TEAL, Pt(1))
add_text(s, Inches(1.0), Inches(6.48), Inches(11), Inches(0.5),
         "SCG-LR consistently Pareto-dominates random_50 across all three effective environments",
         font_size=16, color=ACCENT_TEAL, bold=True)

slide_number(s, 8, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 9: Section Divider — Core Findings
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, SECTION_BG)
add_rect(s, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), ACCENT_GOLD)
add_text(s, Inches(1), Inches(2.3), Inches(5), Inches(0.5),
         "SECTION 02-03", font_size=16, color=ACCENT_GOLD, bold=True)
add_text(s, Inches(1), Inches(2.9), Inches(11), Inches(0.9),
         "Core Empirical Findings", font_size=42, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.7), Inches(10), Inches(0.5),
         "Direction Reversal  +  Signal Replacement  →  Zero prior papers report this", font_size=20, color=ACCENT_GOLD)
slide_number(s, 9, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 10: Finding 1 — Direction Reversal
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Finding 1: Signal-Utility Direction is NOT Fixed", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(10.0), Inches(0.35), "NOVEL FINDING", ACCENT_GOLD)

# Left: explanation
add_shape(s, Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.8), CARD_BG)

add_text(s, Inches(0.6), Inches(1.3), Inches(5), Inches(0.4),
         "token_entropy × Optimizer Utility", font_size=18, color=ACCENT_TEAL, bold=True)

dir_data = [
    ("HotpotQA", "ρ = −0.327 ↘", "High entropy = confused → optimizer useless", ACCENT_RED),
    ("MBPP", "ρ = +0.153 ↗", "High entropy = diverse solutions → optimizer helps", ACCENT_GREEN),
    ("APPS", "ρ = +0.144 ↗", "Same direction as MBPP", ACCENT_GREEN),
    ("WebShop", "ρ = +0.133 ↗", "Same direction as MBPP", ACCENT_GREEN),
]

for i, (env, rho, interp, color) in enumerate(dir_data):
    y = Inches(1.9 + i * 1.1)
    add_shape(s, Inches(0.6), y, Inches(5.4), Inches(0.9), DARK_BG, color, Pt(1))
    add_text(s, Inches(0.8), y + Inches(0.05), Inches(2), Inches(0.35),
             env, font_size=16, color=WHITE, bold=True)
    add_text(s, Inches(2.8), y + Inches(0.05), Inches(3), Inches(0.35),
             rho, font_size=16, color=color, bold=True)
    add_text(s, Inches(0.8), y + Inches(0.4), Inches(5), Inches(0.35),
             interp, font_size=12, color=MED_GRAY)

# Implication box
add_shape(s, Inches(0.4), Inches(5.8), Inches(5.8), Inches(1.0), DARK_BG, ACCENT_GOLD, Pt(1.5))
add_multiline(s, Inches(0.6), Inches(5.85), Inches(5.4), Inches(0.9), [
    ("Implication:", ACCENT_GOLD, True, 14),
    ("Methods assuming fixed direction (CoRefine, SEAG,", WHITE, False, 13),
    ("Think Just Enough) WILL fail on HotpotQA", WHITE, False, 13),
], font_size=13)

# Right: figure
added = add_image_safe(s, "fig13_direction_reversal.png", Inches(6.5), Inches(1.2), width=Inches(6.5))
if not added:
    add_text(s, Inches(7.5), Inches(3.5), Inches(5), Inches(0.5),
             "[fig13_direction_reversal.png]", font_size=16, color=MED_GRAY, align=PP_ALIGN.CENTER)

slide_number(s, 10, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 11: Direction Reversal Robustness
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Robustness: Direction Reversal Survives Artifact Removal", font_size=28, color=WHITE, bold=True)

# Explanation
add_text(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
         "After removing 'finish shortcut' (25.3% of high-U data in HotpotQA), direction reversal still holds:",
         font_size=18, color=LIGHT_GRAY)

# Table
hdr_y = Inches(1.8)
add_rect(s, Inches(1), hdr_y, Inches(11.3), Inches(0.5), ACCENT_TEAL)
for txt, x, w in [("Signal", 1.1, 2.5), ("Full (n=1208) ρ", 3.8, 2.2), ("No finish (n=902) ρ", 6.2, 2.5), ("Change", 8.9, 1.5), ("Still GO?", 10.5, 1.5)]:
    add_text(s, Inches(x), hdr_y + Inches(0.05), Inches(w), Inches(0.4),
             txt, font_size=15, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)

for i, (sig, full, nof, chg, go) in enumerate([
    ("token_entropy", "−0.327", "−0.242", "−26%", "GO"),
    ("evidence_count", "−0.586", "−0.311", "−47%", "GO"),
]):
    y = Inches(2.4 + i * 0.6)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_rect(s, Inches(1), y, Inches(11.3), Inches(0.5), bg)
    add_text(s, Inches(1.2), y + Inches(0.05), Inches(2.3), Inches(0.4), sig, font_size=15, color=WHITE, bold=True)
    add_text(s, Inches(3.8), y + Inches(0.05), Inches(2.2), Inches(0.4), full, font_size=15, color=ACCENT_RED, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(6.2), y + Inches(0.05), Inches(2.5), Inches(0.4), nof, font_size=15, color=ACCENT_RED, align=PP_ALIGN.CENTER, bold=True)
    add_text(s, Inches(8.9), y + Inches(0.05), Inches(1.5), Inches(0.4), chg, font_size=15, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.5), y + Inches(0.05), Inches(1.5), Inches(0.4), go, font_size=15, color=ACCENT_GREEN, align=PP_ALIGN.CENTER, bold=True)

# Figure: wrong direction impact
add_image_safe(s, "fig15_wrong_direction.png", Inches(1), Inches(3.8), width=Inches(11))

slide_number(s, 11, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 12: Finding 2 — Signal Replacement
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Finding 2: Strongest Signal Differs Per Environment", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(10.0), Inches(0.35), "NOVEL FINDING", ACCENT_GOLD)

# Left: signal matrix
add_image_safe(s, "fig14_signal_heatmap.png", Inches(0.4), Inches(1.2), width=Inches(6.2))

# Right: key points
add_shape(s, Inches(6.9), Inches(1.2), Inches(6.2), Inches(5.5), CARD_BG)

findings = [
    ("1. Direction not fixed", "token_entropy: HotpotQA negative, MBPP/APPS/WebShop positive", ACCENT_RED),
    ("2. Strongest signal not fixed", "QA→evidence_count, Code→step_count, Web→state_category", ACCENT_BLUE),
    ("3. Signal TYPE not fixed", "HotpotQA/APPS: continuous; WebShop: categorical", ACCENT_TEAL),
    ("4. step_count direction also reverses!", "MBPP ρ=+0.526 vs APPS ρ=−0.274", ACCENT_GOLD),
]

for i, (title, detail, color) in enumerate(findings):
    y = Inches(1.4 + i * 1.2)
    add_shape(s, Inches(7.1), y, Inches(5.8), Inches(1.0), DARK_BG, color, Pt(1))
    add_text(s, Inches(7.3), y + Inches(0.1), Inches(5.4), Inches(0.35),
             title, font_size=15, color=color, bold=True)
    add_text(s, Inches(7.3), y + Inches(0.5), Inches(5.4), Inches(0.35),
             detail, font_size=12, color=LIGHT_GRAY)

# Conclusion
add_shape(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.7), CARD_BG, ACCENT_GOLD, Pt(1.5))
add_text(s, Inches(0.7), Inches(6.3), Inches(12), Inches(0.5),
         "Pre-selecting a fixed signal + fixed direction (11+ concurrent methods) is SYSTEMATICALLY FRAGILE",
         font_size=16, color=ACCENT_GOLD, bold=True)

slide_number(s, 12, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 13: Wrong-Direction Ablation
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Wrong-Direction Ablation: Catastrophic Consequences", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(10.5), Inches(0.35), "ABLATION", ACCENT_RED)

# Table
hdr_y = Inches(1.1)
add_rect(s, Inches(0.3), hdr_y, Inches(12.7), Inches(0.45), ACCENT_RED)
for txt, x, w in [("Env", 0.4, 1.5), ("Gate", 2.0, 0.9), ("Correct SR", 3.1, 1.5), ("Wrong SR", 4.7, 1.5),
                   ("Δ SR", 6.3, 1.2), ("Wrong RR", 7.6, 1.3), ("Failure Mode", 9.0, 4.0)]:
    add_text(s, Inches(x), hdr_y + Inches(0.02), Inches(w), Inches(0.4),
             txt, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

wd_data = [
    ("HotpotQA", "LR", "96.7%", "62.0%", "−34.5pp", "—", "Active mis-triggering"),
    ("HotpotQA", "MLP", "95.3%", "45.3%", "−51.2pp", "0%", "Complete reversal (zero trigger)"),
    ("HotpotQA", "Prompt", "95.3%", "95.3%", "−1.2pp", "84.5%", "YES-bias masks effect"),
    ("APPS", "LR", "65.0%", "58.5%", "−6.5pp", "0%", "Passive abandonment"),
    ("WebShop", "LR", "43.7%", "7.2%", "−36.5pp", "37.1%", "Complete failure"),
]

for i, (env, gate, csr, wsr, dsr, wrr, fail) in enumerate(wd_data):
    y = Inches(1.65 + i * 0.55)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_rect(s, Inches(0.3), y, Inches(12.7), Inches(0.5), bg)

    vals = [(env, 0.4, 1.5, WHITE), (gate, 2.0, 0.9, LIGHT_GRAY), (csr, 3.1, 1.5, ACCENT_GREEN),
            (wsr, 4.7, 1.5, ACCENT_RED), (dsr, 6.3, 1.2, ACCENT_RED), (wrr, 7.6, 1.3, ACCENT_GOLD),
            (fail, 9.0, 4.0, LIGHT_GRAY)]
    for txt, x, w, c in vals:
        is_bold = "−" in txt and "pp" in txt
        add_text(s, Inches(x), y + Inches(0.08), Inches(w), Inches(0.35),
                 txt, font_size=12, color=c, bold=is_bold, align=PP_ALIGN.CENTER if x < 9 else PP_ALIGN.LEFT)

# Conclusion
add_shape(s, Inches(0.3), Inches(4.5), Inches(12.7), Inches(1.0), CARD_BG, ACCENT_RED, Pt(2))
add_multiline(s, Inches(0.6), Inches(4.6), Inches(12), Inches(0.8), [
    ("Conclusion: Direction is the UNIVERSAL FATAL PREREQUISITE for all learning-based gates", ACCENT_RED, True, 17),
    ("Consistent across environments (HotpotQA, APPS, WebShop) and architectures (LR, MLP, Prompt)", LIGHT_GRAY, False, 14),
], font_size=17)

# Right bottom: figure
add_image_safe(s, "fig6_wrong_direction.png", Inches(1), Inches(5.5), height=Inches(1.7))
add_image_safe(s, "fig15_wrong_direction.png", Inches(7), Inches(5.5), height=Inches(1.7))

slide_number(s, 13, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 14: Section Divider — Phase 5
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, SECTION_BG)
add_rect(s, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), ACCENT_BLUE)
add_text(s, Inches(1), Inches(2.3), Inches(5), Inches(0.5),
         "SECTION 04", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(s, Inches(1), Inches(2.9), Inches(11), Inches(0.9),
         "Phase 5: Method Novelty Enhancement", font_size=42, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.7), Inches(10), Inches(0.5),
         "From ⭐⭐ to ⭐⭐⭐⭐  |  Hand-crafted → Automated  |  Binary → VOC  |  LR → Architecture",
         font_size=20, color=ACCENT_BLUE)
slide_number(s, 14, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 15: Phase 5 Motivation
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Phase 5 Motivation: Why Upgrade?", font_size=28, color=WHITE, bold=True)

# Current assessment cards
dims = [
    ("Empirical Finding", "5/5", "Direction reversal + signal replacement. Zero papers report this.", ACCENT_GREEN, "⭐⭐⭐⭐⭐"),
    ("Experimental Rigor", "4/5", "3 effective envs + 1 boundary + wrong-dir ablation + 3-seed", ACCENT_BLUE, "⭐⭐⭐⭐"),
    ("Method Novelty", "2/5", "LR on 5 hand-crafted features — reviewer may say 'just LR'", ACCENT_RED, "⭐⭐"),
]

for i, (dim, score, desc, color, stars) in enumerate(dims):
    y = Inches(1.2 + i * 1.6)
    card = add_shape(s, Inches(0.5), y, Inches(12.3), Inches(1.4), CARD_BG, color, Pt(1.5))

    add_text(s, Inches(0.8), y + Inches(0.15), Inches(3), Inches(0.4),
             dim, font_size=22, color=color, bold=True)
    add_text(s, Inches(9.5), y + Inches(0.15), Inches(3), Inches(0.4),
             stars, font_size=22, color=color, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(0.8), y + Inches(0.65), Inches(11.5), Inches(0.5),
             desc, font_size=16, color=LIGHT_GRAY)

# Arrow to target
add_shape(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), CARD_BG, ACCENT_GOLD, Pt(2))
add_multiline(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8), [
    ("Phase 5 Goal: Method Novelty ⭐⭐ → ⭐⭐⭐⭐", ACCENT_GOLD, True, 20),
    ("(1) Hand-crafted 5D → LLM hidden state d=2560   (2) Binary → VOC estimation   (3) LR → Architecture innovation", LIGHT_GRAY, False, 14),
], font_size=20)

slide_number(s, 15, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 16: Route A — Cascaded Multi-Fidelity Gate
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
         "Route A: Cascaded Multi-Fidelity Gate", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(9.5), Inches(0.35), "PLAN A", ACCENT_BLUE)

# Cascade diagram
levels = [
    ("L0: Cheap Signal Gate", "~60-70% steps", "Confident → decide (zero overhead)", ACCENT_GREEN, Inches(1.2)),
    ("L1: Hidden-State VOC Probe", "~20-30% steps", "Uncertain from L0 → deeper analysis", ACCENT_BLUE, Inches(3.0)),
    ("L2: Trial Rollout", "~5-10% steps", "Still uncertain → direct VOC observation", ACCENT_GOLD, Inches(4.8)),
]

for name, pct, desc, color, y in levels:
    # Main box
    w = Inches(7)
    box = add_shape(s, Inches(0.8), y, w, Inches(1.4), CARD_BG, color, Pt(2))
    add_text(s, Inches(1.0), y + Inches(0.1), Inches(5), Inches(0.4),
             name, font_size=20, color=color, bold=True)
    add_text(s, Inches(1.0), y + Inches(0.55), Inches(6.5), Inches(0.4),
             desc, font_size=14, color=LIGHT_GRAY)
    # Percentage badge
    tag_w = Inches(1.8)
    tag = add_shape(s, Inches(0.8) + w - tag_w - Inches(0.15), y + Inches(0.15), tag_w, Inches(0.35), color)
    tag.text_frame.paragraphs[0].text = pct
    tag.text_frame.paragraphs[0].font.size = Pt(12)
    tag.text_frame.paragraphs[0].font.color.rgb = WHITE
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Arrow indicators between levels
add_text(s, Inches(3.5), Inches(2.65), Inches(2), Inches(0.35),
         "↓ uncertain", font_size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(3.5), Inches(4.45), Inches(2), Inches(0.35),
         "↓ uncertain", font_size=14, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Right: pros/cons
add_shape(s, Inches(8.3), Inches(1.2), Inches(4.8), Inches(5.4), CARD_BG)

add_text(s, Inches(8.5), Inches(1.3), Inches(4), Inches(0.4),
         "Innovation", font_size=18, color=ACCENT_TEAL, bold=True)
add_multiline(s, Inches(8.5), Inches(1.8), Inches(4.4), Inches(1.5), [
    ("Cascade structure", ACCENT_GREEN, False, 14),
    ("Uncertainty-driven escalation", ACCENT_GREEN, False, 14),
    ("VOC estimation at each level", ACCENT_GREEN, False, 14),
], font_size=14, bullet=True, line_spacing=1.4)

add_text(s, Inches(8.5), Inches(3.3), Inches(4), Inches(0.4),
         "Advantages", font_size=18, color=ACCENT_BLUE, bold=True)
add_multiline(s, Inches(8.5), Inches(3.8), Inches(4.4), Inches(1.0), [
    ("Practical (cost-controllable)", LIGHT_GRAY, False, 14),
    ("Interpretable (per-level stats)", LIGHT_GRAY, False, 14),
], font_size=14, bullet=True, line_spacing=1.4)

add_text(s, Inches(8.5), Inches(4.8), Inches(4), Inches(0.4),
         "Risk", font_size=18, color=ACCENT_RED, bold=True)
add_text(s, Inches(8.7), Inches(5.3), Inches(4.2), Inches(0.4),
         "Cascade may not outperform single-layer probe", font_size=14, color=ACCENT_RED)

# Bottom: novelty rating
add_shape(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5), CARD_BG, ACCENT_BLUE, Pt(1))
add_text(s, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.4),
         "Method Novelty: ⭐⭐⭐  |  Estimated Time: ~2 weeks  |  NeurIPS Acceptance: 70-80%",
         font_size=15, color=ACCENT_BLUE, bold=True)

slide_number(s, 16, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 17: Route B — ICGNet
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
         "Route B: In-Context Gating Network (ICGNet)", font_size=28, color=WHITE, bold=True)
add_tag(s, Inches(9.5), Inches(0.35), "PLAN B", ACCENT_GOLD)

# Architecture diagram
arch_steps = [
    ("1. State Projector", "h ∈ R²⁵⁶⁰ → R¹²⁸", ACCENT_TEAL),
    ("2. Utility Encoder", "U ∈ R → R¹²⁸", ACCENT_BLUE),
    ("3. Self-Attention", "Over calibration context → discover direction & importance", ACCENT_GOLD),
    ("4. Cross-Attention", "Query → context → VOC prediction", ACCENT_GREEN),
]

for i, (name, desc, color) in enumerate(arch_steps):
    y = Inches(1.2 + i * 1.2)
    box = add_shape(s, Inches(0.6), y, Inches(6.5), Inches(1.0), CARD_BG, color, Pt(1.5))
    add_text(s, Inches(0.8), y + Inches(0.1), Inches(3.5), Inches(0.35),
             name, font_size=18, color=color, bold=True)
    add_text(s, Inches(0.8), y + Inches(0.5), Inches(6), Inches(0.35),
             desc, font_size=14, color=LIGHT_GRAY)

# Arrows
for i in range(3):
    y = Inches(2.25 + i * 1.2)
    add_text(s, Inches(3.5), y, Inches(1), Inches(0.3),
             "↓", font_size=18, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Key difference box
add_shape(s, Inches(7.5), Inches(1.2), Inches(5.5), Inches(2.0), CARD_BG, ACCENT_GOLD, Pt(2))
add_text(s, Inches(7.7), Inches(1.3), Inches(5), Inches(0.4),
         "Key Innovation", font_size=20, color=ACCENT_GOLD, bold=True)
add_multiline(s, Inches(7.7), Inches(1.8), Inches(5.1), Inches(1.2), [
    ("Traditional: train(data)→fixed θ→f(x;θ)", ACCENT_RED, False, 15),
    ("ICGNet: f(x, calib_data)→decision", ACCENT_GREEN, True, 15),
    ("One forward pass = 'learn' + decide", WHITE, False, 15),
], font_size=15, line_spacing=1.5)

# Pros/cons
add_shape(s, Inches(7.5), Inches(3.5), Inches(5.5), Inches(3.0), CARD_BG)

add_text(s, Inches(7.7), Inches(3.6), Inches(4), Inches(0.4),
         "Advantages", font_size=18, color=ACCENT_GREEN, bold=True)
add_multiline(s, Inches(7.7), Inches(4.1), Inches(5.1), Inches(1.0), [
    ("High academic novelty", LIGHT_GRAY, False, 14),
    ("Zero re-training for new environments", ACCENT_GREEN, True, 14),
    ("Meta-trained across environments", LIGHT_GRAY, False, 14),
], font_size=14, bullet=True, line_spacing=1.3)

add_text(s, Inches(7.7), Inches(5.3), Inches(4), Inches(0.4),
         "Risk", font_size=18, color=ACCENT_RED, bold=True)
add_text(s, Inches(7.9), Inches(5.8), Inches(5), Inches(0.4),
         "Only 3 environments for meta-training may be insufficient", font_size=14, color=ACCENT_RED)

# Bottom: novelty rating
add_shape(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5), CARD_BG, ACCENT_GOLD, Pt(1))
add_text(s, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.4),
         "Method Novelty: ⭐⭐⭐⭐  |  Estimated Time: ~2 weeks  |  NeurIPS Acceptance: 75-85%",
         font_size=15, color=ACCENT_GOLD, bold=True)

slide_number(s, 17, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 18: Phase 5 Comparison Table
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Phase 5 vs Current vs Competing Methods", font_size=28, color=WHITE, bold=True)

# Comparison table
comp_cols = [("Dimension", 0.3, 2.2), ("SCG-LR\n(Current)", 2.6, 2.0), ("Cascade\n(Phase 5A)", 4.7, 2.0),
             ("ICGNet\n(Phase 5B)", 6.8, 2.0), ("CATTS", 8.9, 2.0), ("ARPO", 11.0, 2.0)]

hdr_y = Inches(1.1)
add_rect(s, Inches(0.2), hdr_y, Inches(12.9), Inches(0.7), ACCENT_TEAL)
for txt, x, w in comp_cols:
    txb = add_text(s, Inches(x), hdr_y + Inches(0.05), Inches(w), Inches(0.6),
                   txt.split("\n")[0], font_size=13, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)
    if "\n" in txt:
        add_text(s, Inches(x), hdr_y + Inches(0.3), Inches(w), Inches(0.35),
                 txt.split("\n")[1], font_size=11, color=DARK_BG, align=PP_ALIGN.CENTER)

comp_data = [
    ("Feature Source", "Hand-crafted 5D", "Hidden state\nd=2560", "Hidden state\nd=2560", "Hand-crafted\nvote signals", "Policy entropy"),
    ("Gate Arch", "LR", "3-level cascade", "Cross-attention", "Single threshold", "RL MLP"),
    ("Cross-env Adapt", "Per-env LR", "Cascade transfer", "Zero re-training", "Per-env grid", "Per-env RL"),
    ("Direction Disc.", "YES", "YES", "YES", "NO", "NO"),
    ("Novelty", "⭐⭐", "⭐⭐⭐", "⭐⭐⭐⭐", "⭐⭐", "⭐⭐⭐"),
]

for i, row in enumerate(comp_data):
    y = Inches(1.9 + i * 0.85)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_rect(s, Inches(0.2), y, Inches(12.9), Inches(0.75), bg)

    # Dimension name
    add_text(s, Inches(0.4), y + Inches(0.15), Inches(2.0), Inches(0.45),
             row[0], font_size=13, color=ACCENT_TEAL, bold=True)

    # Values
    colors = [MED_GRAY, ACCENT_BLUE, ACCENT_GOLD, MED_GRAY, MED_GRAY]
    for j, (val, col_color) in enumerate(zip(row[1:], colors)):
        x = comp_cols[j + 1][1]
        w = comp_cols[j + 1][2]
        txt = val.split("\n")[0] if "\n" in val else val
        c = ACCENT_GREEN if val == "YES" else ACCENT_RED if val == "NO" else col_color
        b = val in ("YES", "NO") or "⭐⭐⭐⭐" in val
        add_text(s, Inches(x), y + Inches(0.1), Inches(w), Inches(0.35),
                 txt, font_size=12, color=c, bold=b, align=PP_ALIGN.CENTER)
        if "\n" in val:
            add_text(s, Inches(x), y + Inches(0.38), Inches(w), Inches(0.35),
                     val.split("\n")[1], font_size=10, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Summary
add_shape(s, Inches(0.2), Inches(6.3), Inches(12.9), Inches(0.6), CARD_BG, ACCENT_GOLD, Pt(1))
add_text(s, Inches(0.5), Inches(6.38), Inches(12.3), Inches(0.45),
         "Both success → NeurIPS acceptance probability rises to 80-85% (from current 65-75%)",
         font_size=16, color=ACCENT_GOLD, bold=True)

slide_number(s, 18, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 19: Section Divider — Paper Storyline
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, SECTION_BG)
add_rect(s, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), ACCENT_TEAL)
add_text(s, Inches(1), Inches(2.3), Inches(5), Inches(0.5),
         "SECTION 05", font_size=16, color=ACCENT_TEAL, bold=True)
add_text(s, Inches(1), Inches(2.9), Inches(11), Inches(0.9),
         "Paper Storyline & Contributions", font_size=42, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(3.7), Inches(10), Inches(0.5),
         "NOT 'we propose a new gate'  →  'we discovered a phenomenon + phenomenon drives method'",
         font_size=18, color=LIGHT_GRAY)
slide_number(s, 19, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 20: Three-Act Structure
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Paper Narrative: Three-Act Structure", font_size=28, color=WHITE, bold=True)

acts = [
    ("Act 1", "Problem Setting", "1 page", ACCENT_TEAL, [
        "Background: Test-time optimizers widely used but costly (5-15×)",
        "Status quo: 11+ concurrent works study 'when to trigger'",
        "Hidden assumption: All methods share UNVERIFIED assumption",
        "— signal-utility direction is fixed",
    ]),
    ("Act 2", "Empirical Finding", "2 pages", ACCENT_GOLD, [
        "Finding 1: Direction reversal across environments",
        "  token_entropy: HotpotQA ρ=−0.327 vs MBPP ρ=+0.153",
        "Finding 2: Signal replacement (identity differs per env)",
        "Quantified cost: Wrong-dir → SR −34.5pp (LR), −51.2pp (MLP)",
    ]),
    ("Act 3", "Method + Validation", "3.5 pages", ACCENT_BLUE, [
        "Findings derive two requirements:",
        "  (1) Auto-discover environment-specific features",
        "  (2) Probe signal-utility direction",
        "Validation: 3 envs + 1 boundary + wrong-dir ablation",
    ]),
]

for i, (act, title, pages, color, bullets) in enumerate(acts):
    x = Inches(0.3 + i * 4.3)
    y = Inches(1.2)

    # Card
    add_shape(s, x, y, Inches(4.1), Inches(5.5), CARD_BG, color, Pt(2))

    # Act label
    tag = add_shape(s, x + Inches(0.15), y + Inches(0.15), Inches(1.0), Inches(0.35), color)
    tag.text_frame.paragraphs[0].text = act
    tag.text_frame.paragraphs[0].font.size = Pt(12)
    tag.text_frame.paragraphs[0].font.color.rgb = WHITE
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Pages tag
    ptag = add_shape(s, x + Inches(4.1) - Inches(1.2), y + Inches(0.15), Inches(1.0), Inches(0.35), DARK_BG)
    ptag.text_frame.paragraphs[0].text = pages
    ptag.text_frame.paragraphs[0].font.size = Pt(11)
    ptag.text_frame.paragraphs[0].font.color.rgb = MED_GRAY
    ptag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    add_text(s, x + Inches(0.15), y + Inches(0.65), Inches(3.8), Inches(0.5),
             title, font_size=22, color=color, bold=True)

    # Bullets
    add_multiline(s, x + Inches(0.15), y + Inches(1.3), Inches(3.8), Inches(3.8),
                  [(b, LIGHT_GRAY if "  " in b else WHITE, "  " not in b, 12 if "  " in b else 13) for b in bullets],
                  font_size=13, line_spacing=1.4)

# Core positioning label
add_text(s, Inches(0.3), Inches(0.35), Inches(6), Inches(0.4),
         "", font_size=1, color=DARK_BG)
add_tag(s, Inches(0.5), Inches(6.85), "CORE FIRE", ACCENT_GOLD)
add_text(s, Inches(2.0), Inches(6.85), Inches(10), Inches(0.4),
         "Act 2 is the CORE contribution — novel empirical finding that no prior paper reports",
         font_size=15, color=ACCENT_GOLD, bold=True)

slide_number(s, 20, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 21: Contribution Ranking & Differentiation
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Contributions & Differentiation Strategy", font_size=28, color=WHITE, bold=True)

# Left: contributions
add_text(s, Inches(0.6), Inches(1.1), Inches(5), Inches(0.4),
         "Contribution Ranking", font_size=20, color=ACCENT_TEAL, bold=True)

contribs = [
    ("C1", "Empirical Finding", "Direction reversal + signal replacement + quantified cost", ACCENT_GOLD, "UNIQUE"),
    ("C2", "Framework", "Direction-Aware Gate: auto feature discovery + direction probe", ACCENT_BLUE, "Phase 5"),
    ("C3", "Validation", "3 envs + 1 boundary + wrong-dir ablation + 3-seed", ACCENT_GREEN, "SOLID"),
    ("C4", "Theory", "VOC ≥ 0 scope limitation + CMDP + dual ascent", ACCENT_TEAL, "BONUS"),
]

for i, (code, name, desc, color, tag_text) in enumerate(contribs):
    y = Inches(1.6 + i * 1.2)
    add_shape(s, Inches(0.5), y, Inches(6.0), Inches(1.0), CARD_BG, color, Pt(1.5))

    # Code
    code_shape = add_shape(s, Inches(0.65), y + Inches(0.1), Inches(0.6), Inches(0.35), color)
    code_shape.text_frame.paragraphs[0].text = code
    code_shape.text_frame.paragraphs[0].font.size = Pt(12)
    code_shape.text_frame.paragraphs[0].font.color.rgb = WHITE
    code_shape.text_frame.paragraphs[0].font.bold = True
    code_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_text(s, Inches(1.4), y + Inches(0.1), Inches(3), Inches(0.35),
             name, font_size=16, color=WHITE, bold=True)
    add_text(s, Inches(1.4), y + Inches(0.5), Inches(5), Inches(0.35),
             desc, font_size=12, color=LIGHT_GRAY)

    # Tag
    tag_color = ACCENT_GOLD if "UNIQUE" in tag_text else ACCENT_BLUE if "Phase" in tag_text else ACCENT_GREEN if "SOLID" in tag_text else MED_GRAY
    add_tag(s, Inches(5.2), y + Inches(0.12), tag_text, tag_color)

# Right: differentiation layers
add_text(s, Inches(7.0), Inches(1.1), Inches(5), Inches(0.4),
         "4-Layer Differentiation", font_size=20, color=ACCENT_TEAL, bold=True)

layers = [
    ("Layer 1: Method", "CROWDED", "Not our selling point", ACCENT_RED, "★★★★★"),
    ("Layer 2: Signal", "Medium", "Probe purpose differs", ACCENT_GOLD, "★★★☆☆"),
    ("Layer 3: Finding", "UNIQUE", "Core selling point!", ACCENT_GREEN, "★☆☆☆☆"),
    ("Layer 4: Framework", "Rare", "Elevates the argument", ACCENT_BLUE, "★★☆☆☆"),
]

for i, (layer, crowd, note, color, stars) in enumerate(layers):
    y = Inches(1.6 + i * 1.2)
    add_shape(s, Inches(6.9), y, Inches(6.0), Inches(1.0), CARD_BG, color, Pt(1.5))

    add_text(s, Inches(7.1), y + Inches(0.1), Inches(3), Inches(0.35),
             layer, font_size=16, color=WHITE, bold=True)
    add_text(s, Inches(10.5), y + Inches(0.1), Inches(2.2), Inches(0.35),
             stars, font_size=14, color=color, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(7.1), y + Inches(0.5), Inches(5.5), Inches(0.35),
             f"{crowd} — {note}", font_size=13, color=LIGHT_GRAY)

# Bottom summary
add_shape(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.6), CARD_BG, ACCENT_GOLD, Pt(1.5))
add_text(s, Inches(0.8), Inches(6.48), Inches(11.5), Inches(0.45),
         "Positioning: Empirical Finding Paper + Method (like Scaling Laws: finding first, method follows)",
         font_size=16, color=ACCENT_GOLD, bold=True)

slide_number(s, 21, TOTAL_SLIDES)


# ─────────────────────────────────────────────────
# SLIDE 22: Next Steps & Timeline
# ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, DARK_BG)
add_text(s, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
         "Next Steps & Timeline", font_size=28, color=WHITE, bold=True)

# Timeline
tasks = [
    ("P0", "Hidden State Infrastructure (Shared)", "2-3 days", ACCENT_RED, "Week 1"),
    ("P0", "Cascaded Gate (5A) Implementation", "5-7 days", ACCENT_RED, "Week 1-2"),
    ("P0", "ICGNet (5B) Implementation", "7-10 days", ACCENT_RED, "Week 1-2"),
    ("P1", "Full Evaluation: 3 envs × 3 seeds", "3-5 days", ACCENT_GOLD, "Week 3"),
    ("P1", "Select 5A/5B winner, update writing guide", "1 day", ACCENT_GOLD, "Week 3"),
    ("P2", "Paper Writing (NeurIPS 9p + appendix)", "2-3 weeks", ACCENT_BLUE, "Week 4-6"),
]

for i, (pri, task, time, color, week) in enumerate(tasks):
    y = Inches(1.2 + i * 0.9)

    # Priority badge
    badge = add_shape(s, Inches(0.6), y + Inches(0.1), Inches(0.6), Inches(0.35), color)
    badge.text_frame.paragraphs[0].text = pri
    badge.text_frame.paragraphs[0].font.size = Pt(12)
    badge.text_frame.paragraphs[0].font.color.rgb = WHITE
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Task bar
    add_shape(s, Inches(1.4), y, Inches(8.5), Inches(0.7), CARD_BG, color, Pt(1))
    add_text(s, Inches(1.6), y + Inches(0.1), Inches(6), Inches(0.4),
             task, font_size=16, color=WHITE, bold=True)
    add_text(s, Inches(1.6), y + Inches(0.4), Inches(3), Inches(0.25),
             time, font_size=12, color=MED_GRAY)

    # Week badge
    wbadge = add_shape(s, Inches(8.5), y + Inches(0.15), Inches(1.2), Inches(0.35), DARK_BG)
    wbadge.text_frame.paragraphs[0].text = week
    wbadge.text_frame.paragraphs[0].font.size = Pt(11)
    wbadge.text_frame.paragraphs[0].font.color.rgb = color
    wbadge.text_frame.paragraphs[0].font.bold = True
    wbadge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Right: summary box
add_shape(s, Inches(10.3), Inches(1.2), Inches(2.7), Inches(5.4), CARD_BG, ACCENT_TEAL, Pt(1.5))
add_text(s, Inches(10.5), Inches(1.3), Inches(2.3), Inches(0.4),
         "Summary", font_size=18, color=ACCENT_TEAL, bold=True)

summary_items = [
    ("Phase 0-4", "DONE", ACCENT_GREEN),
    ("Phase 5", "READY", ACCENT_GOLD),
    ("Total Est.", "~6 weeks", ACCENT_BLUE),
    ("Target", "NeurIPS\n2026", WHITE),
    ("Prob.", "80-85%", ACCENT_GREEN),
]

for i, (label, value, color) in enumerate(summary_items):
    y = Inches(1.8 + i * 0.9)
    add_text(s, Inches(10.5), y, Inches(2.3), Inches(0.3),
             label, font_size=13, color=MED_GRAY)
    add_text(s, Inches(10.5), y + Inches(0.3), Inches(2.3), Inches(0.4),
             value, font_size=18, color=color, bold=True)

# Bottom accent line
add_rect(s, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), ACCENT_TEAL)

slide_number(s, 22, TOTAL_SLIDES)


# ─── Save ───
output_dir = os.path.dirname(__file__)
output_path = os.path.join(output_dir, "drafts", "v1_weekly_report_20260302.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {TOTAL_SLIDES}")
